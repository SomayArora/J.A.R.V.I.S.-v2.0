import os
import random
import re
import logging
import g4f
import g4f.Provider
from flask import Flask, request, jsonify, send_file, render_template
from flask_limiter import Limiter
from pptx import Presentation

app = Flask(__name__)

logging.basicConfig(filename='app.log', level=logging.INFO)

limiter = Limiter(
    app,
    default_limits=["10 per day"],
)

PROMPT_TEMPLATE = """You are a presentation writer AI.

You must follow this format exactly. Do not use Markdown or bullets.

---
#Title: THE TITLE

#Slide: 1
#Header: TABLE OF CONTENTS
#Content: 1. FIRST SLIDE
2. SECOND SLIDE
3. THIRD SLIDE

#Slide: 2
#Header: SLIDE HEADER
#Content: Slide content (short, under 250 characters)

...

#Slide: END
---

Strict rules:
- Use only "#Title:", "#Slide:", "#Header:", and "#Content:".
- Do NOT use markdown (no **bold**, no --- lines, no bullets).
- Every slide must follow the exact same tag format.
- Content under 250 characters per slide.

Begin."""


def create_ppt_text(input_topic):
    """Generate presentation content using the g4f provider."""
    try:
        response = g4f.ChatCompletion.create(
            model="gpt-4o",
            provider=g4f.Provider.Blackbox,  # <-- Free and no API key required
            messages=[
                {"role": "system", "content": PROMPT_TEMPLATE},
                {"role": "user", "content": f"The user wants a presentation about {input_topic}"}
            ],
            stream=True,
        )
    except Exception as e:
        print(f"Error during API call: {e}")
        return ""

    presentation_text = ""
    for message in response:
        if "[DONE]" in str(message):
            continue
        presentation_text += str(message)

    output_path = f'Cache/{input_topic}.txt'
    tagged_text = convert_markdown_to_tagged_format(presentation_text)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(tagged_text)


    return presentation_text

def convert_markdown_to_tagged_format(markdown_text):
    lines = markdown_text.splitlines()
    result = []
    slide_number = 1
    in_slide = False

    for line in lines:
        line = line.strip()

        if line.lower().startswith("### slide"):
            if in_slide:
                result.append("#Slide: END")
            result.append(f"#Slide: {slide_number}")
            slide_number += 1
            in_slide = True
            continue

        if line.lower().startswith("- **title:**"):
            result.append(f"#Title: {line.split(':', 1)[1].strip()}")
            continue

        if line.lower().startswith("- **subtitle:**"):
            result.append(f"#Content: {line.split(':', 1)[1].strip()}")
            continue

        if line.lower().startswith("- **header:**"):
            result.append(f"#Header: {line.split(':', 1)[1].strip()}")
            continue

        if line.lower().startswith("- **content:**"):
            result.append(f"#Content: {line.split(':', 1)[1].strip()}")
            continue

        if line.startswith("- ") or line.startswith("* "):
            result.append(line[2:])
            continue

    if in_slide:
        result.append("#Slide: END")

    return '\n'.join(result)


def create_ppt(text_file, design_number, ppt_name, host_url):
    design_path = f"Designs/Design-{design_number}.pptx"
    print(f"Looking for design file at: {os.path.abspath(design_path)}")

    if not os.path.exists(design_path):
        raise FileNotFoundError(f"Design template not found at {design_path}.")

    prs = Presentation(design_path)
    slide = None
    header = ""
    content = ""

    with open(text_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    for i, line in enumerate(lines):
        line = line.strip()

        if line.startswith("#Title:"):
            title_text = line.replace("#Title:", "").strip()
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
            if slide.shapes.title:
                slide.shapes.title.text = title_text

        elif line.startswith("#Slide:"):
            # Add previous content if exists
            if header or content:
                layout = prs.slide_layouts[1]  # Use consistent layout
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = header
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx != 0:
                        try:
                            placeholder.text = content
                            break
                        except:
                            continue
                header = ""
                content = ""

        elif line.startswith("#Header:"):
            header = line.replace("#Header:", "").strip()

        elif line.startswith("#Content:"):
            content = line.replace("#Content:", "").strip()
            # Handle multi-line content
            j = i + 1
            while j < len(lines):
                next_line = lines[j].strip()
                if next_line.startswith("#"):
                    break
                content += "\n" + next_line
                j += 1

    # Save to file
    output_file_path = f'GeneratedPresentations/{ppt_name}.pptx'
    prs.save(output_file_path)
    abs_path = os.path.abspath(output_file_path)
    os.startfile(abs_path)
    return f"{host_url}{output_file_path}"



@app.route('/GeneratedPresentations/<path:path>')
def send_generated_presentation(path):
    return send_file(f'GeneratedPresentations/{path}', as_attachment=True)


@app.route("/")
def home():
    return render_template("powerpoint.html", charset="utf-8")


@app.route('/generate', methods=['POST'])
@limiter.limit("10 per day")
def generate_presentation():
    topic = request.form.get('topic')
    if not topic:
        return "Please provide a topic.", 400

    ppt_link = get_bot_response(topic, request.host_url)
    return ppt_link


def get_bot_response(msg, host_url="http://localhost:5000/"):
    """Process user input and generate a PowerPoint presentation."""
    user_input = msg.strip()
    last_char = user_input[-1]
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', user_input).rstrip()
    number = 1

    if last_char.isdigit():
        number = int(last_char)
        input_string = user_input[:-2].rstrip()
        print(f"Design Number: {number} selected.")
    else:
        print("No design specified, using default design...")

    os.makedirs('Cache', exist_ok=True)

    ppt_text = create_ppt_text(input_string)
    if not ppt_text:
        return "Error generating presentation text."

    text_file_path = f'Cache/{input_string}.txt'

    if not os.path.exists(text_file_path):
        return f"Error: The file {text_file_path} was not created."

    ppt_link = create_ppt(text_file_path, number, input_string, host_url)
    return str(ppt_link)

# task = input("Task=>")
# if "presentation" in task or "ppt" in task:
#     task=task.replace("jarvis","")
#     task = task.replace("make a presentation", "")
#     task = task.replace("make a ppt", "")
#     task=task.replace("create a presentation","")
#     task=task.replace("create a ppt","")
#     task=task.replace("powerpoint","")
#     task=task.replace("about","")
#     task=task.replace("with design","")
#     topic=task.replace(task[-1],"")
#     task = task.replace(" ","")
#     print(task)
#     print(topic)
#     print("Generating presentation content for you sir. Please hold on.")
#     print("Fetching information to create an engaging presentation.")
#     get_bot_response(task)
#     print(f"Saving the presentation as {topic}.")
#     print(f"Your presentation on {topic} is complete and ready for review sir.")/

get_bot_response("Artificial Intelligence 3")