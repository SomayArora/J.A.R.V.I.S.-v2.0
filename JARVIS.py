import wmi
import smtplib
import random
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
import speedtest
import google.generativeai as genai
from PIL import Image
import random
import sounddevice as sd
import numpy as np
import logging
import g4f
import g4f.Provider
from flask import Flask, request, jsonify, send_file, render_template
from flask_limiter import Limiter
from pptx import Presentation
import re
from io import BytesIO
from ctypes import cast, POINTER
from comtypes import CLSCTX_ALL
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
import pyautogui as gui
import pywhatkit
import datetime
import requests
from playsound import playsound
import os
import psutil
from typing import Union
import sys
import time
import threading
from Alert import Alert
import subprocess
import pywhatkit as pw
import spotipy
import webbrowser
from spotipy.oauth2 import SpotifyOAuth
from rich import print
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options
from os import getcwd
from googletrans import Translator
import pygame
from bs4 import BeautifulSoup

chrome_options = Options()
chrome_options.add_argument("--use-fake-ui-for-media-stream")
chrome_options.add_argument("--headless=new")
chrome_driver_path = f"{getcwd()}\\chromedriver.exe"
service = Service(executable_path=chrome_driver_path)
driver = webdriver.Chrome(service=service, options=chrome_options)

website = "https://allorizenproject1.netlify.app/"
driver.get(website)

Recog_File = f"{getcwd()}\\input.txt"


threshold = 400
clap_detected = False

def detect_clap(indata, frames, time, status):
    global clap_detected
    volume_norm = np.linalg.norm(indata) * 10
    if volume_norm > threshold and not clap_detected:
        print("Clap detected!")
        clap_detected = True

def listen_for_claps():
    with sd.InputStream(callback=detect_clap):
        while True:
            if clap_detected:
                break
            sd.sleep(100)

def main_clap_exe():
    while True:
        print("Listening for claps...")
        listen_for_claps()
        speak("Initialising Jarvis...")
        speak("All Systems will be prepared in a few minutes!!")
        speak("Importing all preferences from home interface and calibrating virtual environment")
        sound_ironman()
        speak("Systems now are fully operational")
        wishMe()
        Jarvis()

def clap_exe():
    print("Listening for claps...")
    listen_for_claps()
    speak("Welcome Back Mister Arora")
    speak("I am At your service sir")
    speak("How can i help you")
    Jarvis()

def clear_file():
    with open(f"{getcwd()}\\input.txt", "w") as file:
        file.truncate(0)

def translate_to_english(text):
    """Translate the given text to English."""
    translator = Translator()
    try:
        translated = translator.translate(text, dest='en')
        return translated.text.lower()
    except:
        return text.lower()

def listen():
    print("CREATED BY SOMAY ARORA")
    output_text = ""
    text = ""

    try:
        start_button = WebDriverWait(driver, 20).until(EC.element_to_be_clickable((By.ID, 'startButton')))
        start_button.click()
        print("Listening...")

        is_second_click = False

        while True:
            output_element = WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.ID, 'output')))
            current_text = output_element.text.strip()

            if "Start Listening" in start_button.text and is_second_click:
                if output_text:
                    is_second_click = False
            elif "Listening..." in start_button.text:
                is_second_click = True

            if current_text != output_text:
                output_text = current_text
                with open(Recog_File, "w") as file:
                    file.write(output_text.lower())
                    print("User:", output_text)


    except KeyboardInterrupt:
        print("Process interrupted by user.")
    except Exception as e:
        print("An error occurred:", e)
    finally:
        driver.quit()

    return text.lower()

def generate_audio(message: str, voice: str = "Matthew"):
    url: str = f"https://api.streamelements.com/kappa/v2/speech?voice={voice}&text={{{message}}}"

    headers = {
        'User-Agent': 'Mozilla/5.0(Maciontosh;intel Mac OS X 10_15_7)AppleWebKit/537.36(KHTML,like Gecoko)Chrome/119.0.0.0 Safari/537.36'}

    try:
        result = requests.get(url=url, headers=headers)
        return result.content
    except:
        return None

def print_animated_message(message):
    if message is None:
        print("Error: Message is None, cannot print.")
        return
    for char in message:
        sys.stdout.write(char)
        sys.stdout.flush()
        time.sleep(0.060)
    print()

def Co_speak(message: str, voice: str = "Matthew", folder: str = "", extension: str = ".mp3") -> Union[None, str]:
    try:
        result_content = generate_audio(message, voice)
        file_path = os.path.join(folder, f"{voice}{extension}")
        with open(file_path, "wb") as file:
            file.write(result_content)
        playsound(file_path)
        os.remove(file_path)
        return None
    except Exception as e:
        print(e)

def speak(text):
    if text is None:
        return

    t1 = threading.Thread(target=Co_speak, args=(text,))
    t2 = threading.Thread(target=print_animated_message, args=(text,))
    t1.start()
    t2.start()
    t1.join()
    t2.join()

def sound_ironman():
    pygame.mixer.init()
    pygame.mixer.music.load('ironman.mp3')
    pygame.mixer.music.play()
    time.sleep(50)
    pygame.mixer.music.stop()

def wishMe():
    hour = datetime.datetime.now().hour

    if hour >= 0 and hour < 12:
        speak("Welcome back Mister Arora")
        speak("Good morning")

    elif hour >= 12 and hour < 17:
        speak("Welcome back Mister Arora")
        speak("Good afternoon")

    else:
        speak("Welcome back Mister Arora")
        speak("Good evening")

def Taskexe():

    def split_and_save_paragraphs(data, filename):
        paragraphs = data.split('\n\n')
        with open(filename, 'w') as file:
            file.write(data)
        data = paragraphs[:2]
        separator = ', '
        joined_string = separator.join(data)
        return joined_string

    def open_App(text):

        try:
            subprocess.run(text)

        except:
            gui.hotkey('win')
            time.sleep(0.5)
            gui.write(text)
            time.sleep(0.5)
            gui.hotkey('enter')

    client_id = "8fd01eeeaf0e4f8d9b352b5d79eec2c9"
    client_secret = "541154d40020419f88d1e877ca770f5b"
    redirect_uri = "http://localhost:8080"

    scope = "user-library-read"

    sp = spotipy.Spotify(auth_manager=SpotifyOAuth(client_id=client_id,
                                                   client_secret=client_secret,
                                                   redirect_uri=redirect_uri,
                                                   scope=scope))

    def play_music_on_spotify(song_name):
        results = sp.search(q=song_name, type='track', limit=1)
        if results['tracks']['items']:
            track = results['tracks']['items'][0]
            print(f"Found: {track['name']} by {track['artists'][0]['name']}")
        else:
            print(f"Song '{song_name}' not found.")

        song = track['external_urls']['spotify']
        webbrowser.open(song)

    def search_playlist(playlist_name):

        results = sp.search(q=playlist_name, type='playlist', limit=1)

        if playlist_name == 'all time':
            #play playlist

        elif results['playlists']['items']:
            playlist = results['playlists']['items'][0]
            print(f"Found Playlist: {playlist['name']}")
            playlist_n = playlist['external_urls']['spotify']
            webbrowser.open(playlist_n)
        else:
            print(f"Playlist '{playlist_name}' not found.")

    def play_music_on_youtube(Song_name):
        pw.playonyt(Song_name)

    def play():
        gui.press("space")

    def search_google(text):
        pywhatkit.search(text)

    def close():
        gui.hotkey('alt', 'f4')

    def search(text):
        gui.press("/")
        time.sleep(0.3)
        gui.write(text)
        gui.hotkey("enter")

    battery = psutil.sensors_battery()

    def battery_Alert():
        while True:
            time.sleep(3)
            percentage = int(battery.percent)
            if percentage == 100:
                t1 = threading.Thread(target=Alert, args=("100% charged",))
                t2 = threading.Thread(target=speak, args=("Sir, The Battery is 100% charged. Please unplug it.",))
                t1.start()
                t2.start()
                t1.join()
                t2.join()
            elif percentage <= 20:
                t1 = threading.Thread(target=Alert, args=("Battery Low",))
                t2 = threading.Thread(target=speak, args=("Sir,Sorry to disturb you but battery is Low now",))
                t1.start()
                t2.start()
                t1.join()
                t2.join()
            elif percentage <= 10:
                t1 = threading.Thread(target=Alert, args=("Battery is too Low",))
                t2 = threading.Thread(target=speak,
                                      args=("Sir,Sorry to disturb you but we are running on very low battery power",))
                t1.start()
                t2.start()
                t1.join()
                t2.join()
            elif percentage <= 5:
                t1 = threading.Thread(target=Alert, args=("Battery is going to die",))
                t2 = threading.Thread(target=speak, args=(
                "Sir,Sorry to disturb you but this is your last chance sir , charge your system now",))
                t1.start()
                t2.start()
                t1.join()
                t2.join()
            time.sleep(10)

    def check_plug():
        print("_____started___")
        battery = psutil.sensors_battery()
        previous_state = battery.power_plugged
        while True:
            battery = psutil.sensors_battery()
            if battery.power_plugged != previous_state:
                if battery.power_plugged:
                    t1 = threading.Thread(target=Alert, args=("Charging **STARTED**",))
                    t2 = threading.Thread(target=speak, args=("Charging Started",))
                    t1.start()
                    t2.start()
                    t1.join()
                    t2.join()
                else:
                    t1 = threading.Thread(target=Alert, args=("Charging **STOPPED**",))
                    t2 = threading.Thread(target=speak, args=("Charging Stopped",))
                    t1.start()
                    t2.start()
                    t1.join()
                    t2.join()

                previous_state = battery.power_plugged

    def check_percentage():
        battery = psutil.sensors_battery()
        percent = int(battery.percent)
        t1 = threading.Thread(target=Alert, args=(f"Battery level at {percent}% sir",))
        t2 = threading.Thread(target=speak, args=(f"Battery at {percent}% power",))
        t1.start()
        t2.start()
        t1.join()
        t2.join()

    def scroll_up():
        gui.press('up', presses=5)

    def scroll_down():
        gui.press('down', presses=5)

    def scroll_to_top():
        gui.hotkey('home')

    def scroll_to_bottom():
        gui.hotkey('end')

    def open_new_tab():
        gui.hotkey('ctrl', 't')

    def close_tab():
        gui.hotkey('ctrl', 'w')

    def open_browser_menu():
        gui.hotkey('alt', 'f')

    def zoom_in():
        gui.hotkey('ctrl', '+')

    def zoom_out():
        gui.hotkey('ctrl', '-')

    def refresh_page():
        gui.hotkey('ctrl', 'r')

    def switch_to_next_tab():
        gui.hotkey('ctrl', 'tab')

    def switch_to_previous_tab():
        gui.hotkey('ctrl', 'shift', 'tab')

    def open_history():
        gui.hotkey('ctrl', 'h')

    def open_bookmarks():
        gui.hotkey('ctrl', 'b')

    def go_back():
        gui.hotkey('alt', 'left')

    def go_forward():
        gui.hotkey('alt', 'right')

    def open_dev_tools():
        gui.hotkey('ctrl', 'shift', 'i')

    def toggle_full_screen():
        gui.hotkey('f11')

    def open_private_window():
        gui.hotkey('ctrl', 'shift', 'n')

    def openweb(webname):

        websites = {
            "youtube": "www.youtube.com",
            "facebook": "www.facebook.com",
            "github": "www.github.com",
            "youtube studio": "studio.youtube.com",
            "twitter": "www.twitter.com",
            "instagram": "www.instagram.com",
            "linkedin": "www.linkedin.com",
            "wikipedia": "www.wikipedia.org",
            "reddit": "www.reddit.com",
            "ova games":"www.ovagames.com",
            "pinterest": "www.pinterest.com",
            "quora": "www.quora.com",
            "tumblr": "www.tumblr.com",
            "flickr": "www.flickr.com",
            "snapchat": "www.snapchat.com",
            "tiktok": "www.tiktok.com",
            "vimeo": "www.vimeo.com",
            "dropbox": "www.dropbox.com",
            "onedrive": "www.onedrive.com",
            "google drive": "drive.google.com",
            "icloud": "www.icloud.com",
            "amazon": "www.amazon.com",
            "ebay": "www.ebay.com",
            "alibaba": "www.alibaba.com",
            "netflix": "www.netflix.com",
            "hulu": "www.hulu.com",
            "disney plus": "www.disneyplus.com",
            "hbo max": "www.hbomax.com",
            "spotify": "www.spotify.com",
            "soundcloud": "www.soundcloud.com",
            "apple music": "www.apple.com/apple-music",
            "pandora": "www.pandora.com",
            "deezer": "www.deezer.com",
            "bandcamp": "www.bandcamp.com",
            "bbc": "www.bbc.com",
            "cnn": "www.cnn.com",
            "nytimes": "www.nytimes.com",
            "the guardian": "www.theguardian.com",
            "forbes": "www.forbes.com",
            "bloomberg": "www.bloomberg.com",
            "reuters": "www.reuters.com",
            "espn": "www.espn.com",
            "fox news": "www.foxnews.com",
            "nbc news": "www.nbcnews.com",
            "cbs news": "www.cbsnews.com",
            "abc news": "www.abcnews.go.com",
            "msnbc": "www.msnbc.com",
            "npr": "www.npr.org",
            "wsj": "www.wsj.com",
            "yahoo news": "news.yahoo.com",
            "buzzfeed": "www.buzzfeed.com",
            "huffpost": "www.huffpost.com",
            "canva": "www.canva.com",
            "chatgpt": "chat.openai.com",
            "slack": "www.slack.com",
            "trello": "www.trello.com",
            "asana": "www.asana.com",
            "zoom": "www.zoom.us",
            "skype": "www.skype.com",
            "microsoft teams": "www.microsoft.com/microsoft-teams",
            "google meet": "meet.google.com",
            "webex": "www.webex.com",
            "jira": "www.atlassian.com/software/jira",
            "notion": "www.notion.so",
            "airtable": "www.airtable.com",
            "monday": "www.monday.com",
            "clickup": "www.clickup.com",
            "dropbox paper": "www.dropbox.com/paper",
            "confluence": "www.atlassian.com/software/confluence",
            "figma": "www.figma.com",
            "adobe xd": "www.adobe.com/products/xd.html",
            "invision": "www.invisionapp.com",
            "microsoft word": "www.microsoft.com/microsoft-365/word",
            "google docs": "docs.google.com",
            "medium": "www.medium.com",
            "wordpress": "www.wordpress.com",
            "wix": "www.wix.com",
            "squarespace": "www.squarespace.com",
            "shopify": "www.shopify.com",
            "bigcommerce": "www.bigcommerce.com",
            "weebly": "www.weebly.com",
            "godaddy": "www.godaddy.com",
            "namecheap": "www.namecheap.com",
            "bluehost": "www.bluehost.com",
            "siteground": "www.siteground.com",
            "hostgator": "www.hostgator.com",
            "dreamhost": "www.dreamhost.com",
            "a2 hosting": "www.a2hosting.com",
            "inmotion hosting": "www.inmotionhosting.com",
            "digitalocean": "www.digitalocean.com",
            "linode": "www.linode.com",
            "aws": "aws.amazon.com",
            "azure": "azure.microsoft.com",
            "google cloud": "cloud.google.com",
            "heroku": "www.heroku.com",
            "gitlab": "www.gitlab.com",
            "bitbucket": "bitbucket.org",
            "codepen": "codepen.io",
            "jsfiddle": "jsfiddle.net",
            "repl.it": "repl.it",
            "stack overflow": "stackoverflow.com",
            "stackoverflow careers": "stackoverflow.com/jobs",
            "glassdoor": "www.glassdoor.com",
            "indeed": "www.indeed.com",
            "linkedin jobs": "www.linkedin.com/jobs",
            "monster": "www.monster.com",
            "simplyhired": "www.simplyhired.com",
            "angel.co": "angel.co",
            "github jobs": "jobs.github.com",
            "ziprecruiter": "www.ziprecruiter.com",
            "careerbuilder": "www.careerbuilder.com",
            "snagajob": "www.snagajob.com",
            "dice": "www.dice.com",
            "jobs": "www.jobs.com",
            "bamboohr": "www.bamboohr.com",
            "workday": "www.workday.com",
            "adp": "www.adp.com",
            "sap successfactors": "www.sap.com/products/hcm.html",
            "oracle hcm": "www.oracle.com/applications/human-capital-management",
            "zenefits": "www.zenefits.com",
            "paycor": "www.paycor.com",
            "paycom": "www.paycom.com",
            "gusto": "www.gusto.com",
            "square": "squareup.com",
            "stripe": "www.stripe.com",
            "paypal": "www.paypal.com",
            "venmo": "www.venmo.com",
            "cash app": "cash.app",
            "robinhood": "www.robinhood.com",
            "etrade": "www.etrade.com",
            "fidelity": "www.fidelity.com",
            "charles schwab": "www.schwab.com",
            "vanguard": "investor.vanguard.com",
            "td ameritrade": "www.tdameritrade.com",
            "coinbase": "www.coinbase.com",
            "binance": "www.binance.com",
            "kraken": "www.kraken.com",
            "blockchain": "www.blockchain.com",
            "gemini": "www.gemini.com",
            "bitfinex": "www.bitfinex.com",
            "bitstamp": "www.bitstamp.net",
            "bittrex": "www.bittrex.com",
            "okex": "www.okex.com",
            "poloniex": "www.poloniex.com",
            "coindesk": "www.coindesk.com",
            "cointelegraph": "www.cointelegraph.com",
            "decrypt": "www.decrypt.co",
            "cryptoslate": "www.cryptoslate.com",
            "cryptonews": "www.cryptonews.com",
            "coinmarketcap": "www.coinmarketcap.com",
            "coingecko": "www.coingecko.com",
            "messari": "www.messari.io",
            "icodrops": "www.icodrops.com",
            "tokenmarket": "www.tokenmarket.net",
            "coinpaprika": "www.coinpaprika.com",
            "cryptocompare": "www.cryptocompare.com",
            "coincheckup": "www.coincheckup.com",
            "cryptobriefing": "www.cryptobriefing.com",
            "blockonomi": "www.blockonomi.com",
            "coininsider": "www.coininsider.com",
            "newsbtc": "www.newsbtc.com",
            "bitcoin.com": "www.bitcoin.com",
            "ethereum.org": "www.ethereum.org",
            "litecoin.com": "www.litecoin.com",
            "ripple.com": "www.ripple.com",
            "cardano.org": "www.cardano.org",
            "stellarlumens.com": "www.stellarlumens.com",
            "tezos.com": "www.tezos.com",
            "eos.io": "www.eos.io",
            "neo.org": "www.neo.org",
            "iota.org": "www.iota.org",
            "monero.org": "www.monero.org",
            "zcash.org": "www.zcash.org",
            "dash.org": "www.dash.org",
            "dogecoin.com": "www.dogecoin.com",
            "gpt": "www.chatgpt.com/"
        }

        webname_lower = webname.lower()
        webname_lower = webname_lower.replace(" ", "")

        if webname_lower in websites:
            webbrowser.open(f"http://{websites[webname_lower]}")
            print(f"Opening {webname}...")
        else:
            web = 'https://www.' + webname_lower + '.com'
            webbrowser.open(web)
            print(f"Opening {webname}...")

    def volume_up():
        gui.press('up')

    def volume_down():
        gui.press('down')

    def seek_forward():
        gui.press('right')

    def seek_backward():
        gui.press('left')

    def seek_forward_10s():
        gui.press('l')

    def seek_backward_10s():
        gui.press('j')

    def seek_backward_frame():
        gui.press(',')

    def seek_forward_frame():
        gui.press('.')

    def seek_to_beginning():
        gui.press('home')

    def seek_to_end():
        gui.press('end')

    def seek_to_previous_chapter():
        gui.hotkey('ctrl', 'left')

    def seek_to_next_chapter():
        gui.hotkey('ctrl', 'right')

    def decrease_playback_speed():
        gui.hotkey('shift', ',')

    def increase_playback_speed():
        gui.hotkey('shift', '.')

    def move_to_next_video():
        gui.hotkey('shift', 'n')

    def move_to_previous_video():
        gui.hotkey('shift', 'p')

    online_dlg = [
        "Sir, I am online and ready",
        "Sir, I am online",
        "I am online, sir",
        "Sir, my status is online",
        "I am online, sir, happy to assist you",
        "Sir, I am now online",
        "I am online and ready to help, sir",
        "Sir, I have gone online",
        "Sir, I am online, awaiting your command",
        "Sir, I am online and at your service",
        "Online and ready, sir",
        "I am online and here to assist, sir",
        "Sir, I am now online and ready",
        "Sir, I am online and ready to serve",
        "Sir, I am online and operational",
        "Sir, I have gone online and am ready to assist",
        "Sir, I am online and available",
        "Online and at your service, sir",
        "Sir, I am online and prepared to help",
        "Sir, I am online and standing by",
        "Sir, I am online, how can I assist?",
        "Sir, I am online and awaiting your instructions",
        "Sir, I am now online, how can I help?",
        "Sir, I am online and ready for your commands",
        "Sir, I am online and ready to take your orders",
        "Sir, I am online and ready for duty",
        "Online and ready to assist, sir",
        "Sir, I am online, what do you need?",
        "Sir, I am online, how can I serve you?",
        "Sir, I am online and at your disposal",
        "Sir, I am online, what can I do for you?",
        "Sir, I am online, ready and waiting",
        "Sir, I am online and ready to go",
        "Sir, I am online and ready to support you",
        "Sir, I am online, your command is my wish",
        "Sir, I am online and ready for action",
        "Sir, I am online, ready to serve",
        "Sir, I am online and here to help",
        "Sir, I am online, how can I be of assistance?",
        "Sir, I am online, how may I assist you?",
        "Sir, I am online and at your service, awaiting commands",
        "Sir, I am online and ready to execute your commands",
        "Sir, I am online and ready to perform tasks",
        "Sir, I am online and ready to assist with any tasks",
        "Sir, I am online and ready to aid you",
        "Sir, I am online and ready to handle your requests",
        "Sir, I am online and here to serve you",
        "Sir, I am online and at your command",
        "Sir, I am online and awaiting your requests",
        "Sir, I am online and ready for instructions",
        "Sir, I am online and prepared to assist with any task",
        "Sir, I am online and ready to follow your orders",
        "Sir, I am online and prepared to serve",
        "Sir, I am online, ready to assist",
        "Sir, I am online and ready to support",
        "Sir, I am online, here to help",
        "Sir, I am online and ready to execute orders",
        "Sir, I am online and ready for your instructions",
        "Sir, I am online and ready to provide assistance",
        "Sir, I am online and awaiting your commands",
        "Sir, I am online and here to aid you",
        "Sir, I am online and ready for your tasks",
        "Sir, I am online and ready for any commands",
        "Sir, I am online and ready to work",
        "Sir, I am online and ready to be of service",
        "Sir, I am online and at your disposal, ready to help",
        "Sir, I am online and ready to take on tasks",
        "Sir, I am online and prepared to assist in any way",
        "Sir, I am online and ready to fulfill your requests",
        "Sir, I am online and ready to help you",
        "Sir, I am online and prepared for your instructions",
        "Sir, I am online and ready to follow your lead",
        "Sir, I am online and at your beck and call",
        "Sir, I am online and ready to serve you",
        "Sir, I am online and ready to assist with any requests",
        "Sir, I am online and ready to follow your commands",
        "Sir, I am online and here to assist with anything",
        "Sir, I am online and ready for any task",
        "Sir, I am online and ready for your guidance",
        "Sir, I am online and ready to assist you in any way",
        "Sir, I am online and ready to perform",
        "Sir, I am online and at your service, ready to help",
        "Sir, I am online and ready for any instructions",
        "Sir, I am online and ready to assist with all tasks",
        "Sir, I am online and prepared for any requests",
        "Sir, I am online and ready to take on any task",
        "Sir, I am online and ready for your directions",
        "Sir, I am online and ready to assist with whatever you need",
        "Sir, I am online and ready to serve your needs",
        "Sir, I am online and ready to execute your tasks",
        "Sir, I am online and ready to be at your service",
        "Sir, I am online and ready to assist in any tasks",
        "Sir, I am online and ready to perform any duty",
        "Sir, I am online and prepared to assist you",
        "Sir, I am online and ready for your requests",
        "Sir, I am online and ready for any duty",
        "Sir, I am online and ready to fulfill any request",
    ]
    offline_dlg = [
        "Sir, I am offline. Please connect to the WiFi or internet",
        "Sir, I am currently offline. Please check the internet connection",
        "I am offline, sir. Kindly connect me to the internet",
        "Sir, my status is offline. Please ensure I am connected to the network",
        "I am offline, sir. Please establish an internet connection",
        "Sir, I am offline. Please connect me to the internet",
        "Sir, I am offline. Kindly check the network connection",
        "Sir, I am offline. Please reconnect me to the WiFi",
        "Sir, I am offline. A network connection is required",
        "Sir, I am offline. Please connect to the internet",
        "Offline, sir. Please check the WiFi connection",
        "Sir, I am offline. Internet connection needed",
        "Sir, I am offline. Please verify the internet connection",
        "Sir, I am offline. Network connection required",
        "Sir, I am offline. Please reconnect to the internet",
        "Sir, I am offline. WiFi connection needed",
        "Sir, I am offline. Please ensure network connectivity",
        "Offline status, sir. Please connect to the WiFi",
        "Sir, I am offline. Please verify the network connection",
        "Sir, I am offline. Internet access required",
        "Sir, I am offline. Please check the internet connection",
        "Sir, I am offline. Please ensure a stable network connection",
        "Sir, I am offline. Please reconnect to the network",
        "Sir, I am offline. A WiFi connection is needed",
        "Sir, I am offline. Internet connection is necessary",
        "Sir, I am offline. Please check WiFi or internet",
        "Sir, I am offline. Network access required",
        "Sir, I am offline. Please establish a connection to the internet",
        "Sir, I am offline. Internet connection needed",
        "Sir, I am offline. Please connect me to the network",
        "Sir, I am offline. Please ensure internet connectivity",
        "Sir, I am offline. Network connection needed",
        "Sir, I am offline. Please verify internet access",
        "Sir, I am offline. Please connect to WiFi",
        "Sir, I am offline. Please ensure a network connection",
        "Sir, I am offline. Please connect me to the WiFi",
        "Sir, I am offline. Please reconnect to the internet",
        "Sir, I am offline. Please check the network connection",
        "Sir, I am offline. WiFi connection required",
        "Sir, I am offline. Network connectivity is needed",
        "Sir, I am offline. Internet access required",
        "Sir, I am offline. Please ensure a stable internet connection",
        "Sir, I am offline. Please connect to the network",
        "Sir, I am offline. Internet connection required",
        "Sir, I am offline. Please establish a network connection",
        "Sir, I am offline. Please reconnect to the WiFi",
        "Sir, I am offline. Please ensure a stable WiFi connection",
        "Sir, I am offline. Network access needed",
        "Sir, I am offline. Please connect to the internet",
        "Sir, I am offline. WiFi access required",
        "Sir, I am offline. Please ensure network access",
        "Sir, I am offline. Please check the WiFi connection",
        "Sir, I am offline. Please ensure an internet connection",
        "Sir, I am offline. Network connection required",
        "Sir, I am offline. Internet connection needed",
        "Sir, I am offline. Please connect me to the network",
        "Sir, I am offline. Please ensure internet access",
        "Sir, I am offline. Network connectivity required",
        "Sir, I am offline. Please check the network connection",
        "Sir, I am offline. Please reconnect to the WiFi",
        "Sir, I am offline. Internet connectivity required",
        "Sir, I am offline. Network access needed",
        "Sir, I am offline. Please verify the internet connection",
        "Sir, I am offline. WiFi connection needed",
        "Sir, I am offline. Please check the internet access",
        "Sir, I am offline. Please ensure a WiFi connection",
        "Sir, I am offline. Please reconnect to the network",
        "Sir, I am offline. Internet access needed",
        "Sir, I am offline. Network connection is required",
        "Sir, I am offline. Please connect to the WiFi",
        "Sir, I am offline. Please ensure internet connectivity",
        "Sir, I am offline. Network access required",
        "Sir, I am offline. Please verify network connection",
        "Sir, I am offline. Please check the WiFi connection",
        "Sir, I am offline. Internet connection required",
        "Sir, I am offline. Network connection needed",
        "Sir, I am offline. Please connect to the internet",
        "Sir, I am offline. Please ensure network access",
        "Sir, I am offline. Please check the network connectivity",
        "Sir, I am offline. Please reconnect to the internet",
        "Sir, I am offline. Internet access required",
        "Sir, I am offline. Please ensure a network connection",
        "Sir, I am offline. Please verify internet connection",
        "Sir, I am offline. Please check the WiFi or internet connection",
        "Sir, I am offline. Network connection required",
        "Sir, I am offline. Please reconnect to the network",
        "Sir, I am offline. Please connect to WiFi or internet",
        "Sir, I am offline. Please ensure WiFi access",
        "Sir, I am offline. Internet connection is needed",
        "Sir, I am offline. Please check WiFi or internet access",
        "Sir, I am offline. Please establish network connectivity",
    ]

    def get_ram_info():
        ram = psutil.virtual_memory()
        total_ram = ram.total / (1024 ** 3)
        available_ram = ram.available / (1024 ** 3)
        return ("Assessing total and available RAM.\n"
                f"Sir,Your system is equipped with a total of {total_ram:.2f} gigabytes of RAM.\n"
                f"Currently, {available_ram:.2f} gigabytes of RAM are available for use.")

    def get_storage_info(drive_letter):
        drive_letter = drive_letter.upper()
        partition_info = psutil.disk_partitions()
        for partition in partition_info:
            if partition.device.startswith(drive_letter):
                usage = psutil.disk_usage(partition.mountpoint)
                total = usage.total / (1024 ** 3)
                used = usage.used / (1024 ** 3)
                free = usage.free / (1024 ** 3)
                return (f"Calculating total, used, and available space in Drive {drive_letter}\n"
                        f"Sir, Drive {drive_letter} in your system has a total storage capacity of {total:.2f} gigabytes.\n"
                        f"Currently, {used:.2f} gigabytes are being utilized.\n"
                        f"This leaves you with {free:.2f} gigabytes of available space.")
        return f"Drive {drive_letter} not found Sir."

    def get_brightness_windows():
        try:
            w = wmi.WMI(namespace='wmi')
            brightness_methods = w.WmiMonitorBrightness()
            brightness_percentage = brightness_methods[0].CurrentBrightness
            return brightness_percentage
        except Exception as e:
            return f"Error: {e}"

    def check_br_percentage():
        brightness = get_brightness_windows()
        speak(f"Sir, the Current Brightness is {brightness}%")

    def SpeedTest():

        speak("Initiating internet speed test sir. This will only take a few seconds.")
        speed = speedtest.Speedtest()
        speed.get_best_server()
        downloading = speed.download()
        correctDown = int(downloading/800000)
        uploading = speed.upload()
        correctUpload = int(uploading/800000)
        ping = speed.results.ping

        speak(f"Sir,The Downloading speed is {correctDown} mbp s and The Uploading speed is {correctUpload} mbp s and the ping is {ping}")

    def get_running_apps_windows():
        try:
            # Get a list of running processes
            processes = [proc.name() for proc in psutil.process_iter(['name'])]
            return list(set(processes))  # Remove duplicates
        except Exception as e:
            return f"Error: {e}"

    def check_running_app():
        running_apps = get_running_apps_windows()
        speak("Sir the Running Apps are below")
        for app in running_apps:
            print(app)

    def get_file_extension(text):
        if "python file" in text:
            ex = ".py"
        elif "java file" in text:
            ex = ".java"
        elif "text file" in text:
            ex = ".txt"
        elif "html file" in text:
            ex = ".html"
        elif "css file" in text:
            ex = ".css"
        elif "javascript file" in text:
            ex = ".js"
        elif "json file" in text:
            ex = ".json"
        elif "xml file" in text:
            ex = ".xml"
        elif "csv file" in text:
            ex = ".csv"
        elif "markdown file" in text:
            ex = ".md"
        elif "yaml file" in text:
            ex = ".yaml"
        elif "image file" in text:
            ex = ".jpg"
        elif "video file" in text:
            ex = ".mp4"
        elif "audio file" in text:
            ex = ".mp3"
        elif "pdf file" in text:
            ex = ".pdf"
        elif "word file" in text:
            ex = ".docx"
        elif "excel file" in text:
            ex = ".xlsx"
        elif "powerpoint file" in text:
            ex = ".pptx"
        elif "zip file" in text:
            ex = ".zip"
        elif "tar file" in text:
            ex = ".tar"
        else:
            ex = ""  # Default case if no match found
        return ex

    def update_text(text):
        if "python file" in text:
            text = text.replace("python file", "")
        elif "java file" in text:
            text = text.replace("java file", "")
        elif "text file" in text:
            text = text.replace("text file", "")
        elif "html file" in text:
            text = text.replace("html file", "")
        elif "css file" in text:
            text = text.replace("css file", "")
        elif "javascript file" in text:
            text = text.replace("javascript file", "")
        elif "json file" in text:
            text = text.replace("json file", "")
        elif "xml file" in text:
            text = text.replace("xml file", "")
        elif "csv file" in text:
            text = text.replace("csv file", "")
        elif "markdown file" in text:
            text = text.replace("markdown file", "")
        elif "yaml file" in text:
            text = text.replace("yaml file", "")
        elif "image file" in text:
            text = text.replace("image file", "")
        elif "video file" in text:
            text = text.replace("video file", "")
        elif "audio file" in text:
            text = text.replace("audio file", "")
        elif "pdf file" in text:
            text = text.replace("pdf file", "")
        elif "word file" in text:
            text = text.replace("word file", "")
        elif "excel file" in text:
            text = text.replace("excel file", "")
        elif "powerpoint file" in text:
            text = text.replace("powerpoint file", "")
        elif "zip file" in text:
            text = text.replace("zip file", "")
        elif "tar file" in text:
            text = text.replace("tar file", "")
        else:
            pass
        return text

    def create_file(text):
        selected_ex = get_file_extension(text)
        text = update_text(text)
        if "named" in text or "with name" in text:
            text = text.replace("named", "")
            text = text.replace("with name", "")
            text = text.replace("create a file", "")
            text = text.replace("create a project","")
            text = text.strip()
            with open(f"{text}{selected_ex}", "w"):
                pass
        else:
            with open(f"demo{selected_ex}", "w"):
                pass

    def find_my_ip():
        ip_address = requests.get('https://api64.ipify.org?format=json').json()
        return ip_address["ip"]

    def set_brightness_windows(percentage):
        try:
            w = wmi.WMI(namespace='wmi')
            brightness_methods = w.WmiMonitorBrightnessMethods()[0]
            brightness_methods.WmiSetBrightness(int(percentage), 0)
            speak(f"Brightness set to {percentage}%")
        except Exception as e:
            speak(f"Error: {e}")

    def get_volume_windows():
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        current_volume = volume.GetMasterVolumeLevelScalar() * 100
        speak(f"the device is running on {int(round(current_volume, 2))}  % volume level")

    def set_volume_windows(percentage):
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        volume.SetMasterVolumeLevelScalar(percentage / 100, None)
        speak(f"Volume set to {percentage}%")

    def generate_image(text):

        url = 'https://api.airforce/v1/imagine2'
        params = {'prompt': text}
        response = requests.get(url, params=params)
        if response.status_code == 200:
            image = Image.open(BytesIO(response.content))
            name = text.replace("seed=", "").replace(" ","_")
            image.save(f'{name}.png')
            print(f'Image saved as {name}.png')
            image.show()


        else:
            print(f'Failed to retrieve image. Status code: {response.status_code}')

    def alarm():
        speak("Ok Sir, Enter the time!!")
        time = input("<=Enter The Time=>")

        while True:
            Time_At = datetime.datetime.now()
            now = Time_At.strftime("%H:%M:%S")

            if now == time:
                speak("Time To Wake Up Sir!!")
                playsound('ironman.mp3')
                speak("Alarm Closed!!")

            elif now > time:
                break





    #while True:

        text  = listen()

    def get_weather_by_address(address):
        search_url = f"https://www.google.com/search?q=weather+{address.replace(' ', '+')}"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(search_url, headers=headers)
        if response.status_code == 200:
            soup = BeautifulSoup(response.text, 'html.parser')

            location = soup.find("div", attrs={"id": "wob_loc"}).text
            time = soup.find("div", attrs={"id": "wob_dts"}).text
            weather = soup.find("span", attrs={"id": "wob_dc"}).text
            temp = soup.find("span", attrs={"id": "wob_tm"}).text

            weather_report = (f"Weather: {weather}\n"
                              f"Temperature: {temp}°C")

            speak(f"Sir the weather is {weather} and the temperature is {temp} degrees celcius")
        else:
            return "Error retrieving weather data."

    def send_email(body, receiver_email, subject=""):

        sender_email = "somayarora8008@gmail.com"
        sender_password = "ggsnzpnuclyldxpb"
        display_name = "J.A.R.V.I.S."

        msg = MIMEMultipart()
        msg['From'] = f"{display_name} <{sender_email}>"
        msg['To'] = receiver_email
        msg['Subject'] = subject

        msg.attach(MIMEText(body, 'plain'))

        try:
            server = smtplib.SMTP('smtp.gmail.com', 587)
            server.starttls()

            server.login(sender_email, sender_password)

            server.sendmail(sender_email, receiver_email, msg.as_string())
            print("Email sent successfully!")

        except Exception as e:
            print(f"Failed to send email: {e}")

        finally:
            server.quit()

    def get_date():

        now = datetime.datetime.now()
        date_today = now.strftime("%B %d, %Y")
        speak(f"Sir,Today is {date_today}.")

    def get_time():
        now = datetime.datetime.now()
        hour_now = now.strftime("%I")
        minute_now = now.strftime("%M")
        am_pm = now.strftime("%p")

        speak(f"Sir, the time is {hour_now}:{minute_now} {am_pm}.")

    app = Flask(__name__)

    logging.basicConfig(filename='app.log', level=logging.INFO)

    limiter = Limiter(
        app,
        default_limits=["10 per day"],
    )

    PROMPT_TEMPLATE = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.
    Notice
    -You do all the presentation text for the user.
    -You write the texts no longer than 250 characters!
    -You make very short titles!
    -You make the presentation easy to understand.
    -The presentation has a table of contents.
    -The presentation has a summary.
    -At least 8 slides.

    Example! - Stick to this formatting exactly!
    #Title: TITLE OF THE PRESENTATION

    #Slide: 1
    #Header: table of contents
    #Content: 1. CONTENT OF THIS POWERPOINT
    2. CONTENTS OF THIS POWERPOINT
    3. CONTENT OF THIS POWERPOINT
    ...

    #Slide: 2
    #Header: TITLE OF SLIDE
    #Content: CONTENT OF THE SLIDE

    #Slide: 3
    #Header: TITLE OF SLIDE
    #Content: CONTENT OF THE SLIDE

    #Slide: 4
    #Header: TITLE OF SLIDE
    #Content: CONTENT OF THE SLIDE

    #Slide: 5
    #Headers: summary
    #Content: CONTENT OF THE SUMMARY

    #Slide: END
    """

    def create_ppt_text(input_topic):
        """Generate presentation content using the g4f provider."""
        try:
            response = g4f.ChatCompletion.create(
                model="gpt-3.5-turbo",
                provider=g4f.Provider.ChatGpt,
                messages=[
                    {"role": "system", "content": PROMPT_TEMPLATE},
                    {"role": "user", "content": f"The user wants a presentation about {input_topic}"}
                ],
                stream=True,
            )
        except Exception as e:
            print(f"Error during API call: {e}")
            return ""

        presentation_text = ""
        for message in response:
            if "[DONE]" in str(message):
                continue
            presentation_text += str(message)

        output_path = f'Cache/{input_topic}.txt'
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(presentation_text)

        return presentation_text

    def create_ppt(text_file, design_number, ppt_name, host_url):
        """Create a PowerPoint presentation using the specified design template and open it."""
        design_path = f"Designs/Design-{design_number}.pptx"
        print(f"Looking for design file at: {os.path.abspath(design_path)}")

        if not os.path.exists(design_path):
            raise FileNotFoundError(f"Design template not found at {design_path}. Please check the file path.")

        prs = Presentation(design_path)
        slide_count = 0
        header = ""
        content = ""
        last_slide_layout_index = -1

        with open(text_file, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if line.startswith('#Title:'):
                    header = line.replace('#Title:', '').strip()
                    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
                    slide.shapes.title.text = header
                    continue

                elif line.startswith('#Slide:'):
                    if slide_count > 0:
                        slide = prs.slides.add_slide(prs.slide_layouts[last_slide_layout_index])
                        slide.shapes.title.text = header
                        slide.shapes.placeholders[1].text = content

                    content = ""
                    slide_count += 1
                    last_slide_layout_index = random.choice([1, 7, 8])
                    continue

                elif line.startswith('#Header:'):
                    header = line.replace('#Header:', '').strip()
                    continue

                elif line.startswith('#Content:'):
                    content = line.replace('#Content:', '').strip()
                    next_line = f.readline().strip()
                    while next_line and not next_line.startswith('#'):
                        content += '\n' + next_line
                        next_line = f.readline().strip()
                    continue

        output_file_path = f'GeneratedPresentations/{ppt_name}.pptx'
        prs.save(output_file_path)

        abs_path = os.path.abspath(output_file_path)
        os.startfile(abs_path)

        return f"{host_url}{output_file_path}"

    @app.route('/GeneratedPresentations/<path:path>')
    def send_generated_presentation(path):
        return send_file(f'GeneratedPresentations/{path}', as_attachment=True)

    @app.route("/")
    def home():
        return render_template("powerpoint.html", charset="utf-8")

    @app.route('/generate', methods=['POST'])
    @limiter.limit("10 per day")
    def generate_presentation():
        topic = request.form.get('topic')
        if not topic:
            return "Please provide a topic.", 400

        ppt_link = get_bot_response(topic, request.host_url)
        return ppt_link

    def get_bot_response(msg, host_url="http://localhost:5000/"):
        """Process user input and generate a PowerPoint presentation."""
        user_input = msg.strip()
        last_char = user_input[-1]
        input_string = re.sub(r'[^\w\s.\-\(\)]', '', user_input).rstrip()
        number = 1

        if last_char.isdigit():
            number = int(last_char)
            input_string = user_input[:-2].rstrip()
            print(f"Design Number: {number} selected.")
        else:
            print("No design specified, using default design...")

        os.makedirs('Cache', exist_ok=True)

        ppt_text = create_ppt_text(input_string)
        if not ppt_text:
            return "Error generating presentation text."

        text_file_path = f'Cache/{input_string}.txt'

        if not os.path.exists(text_file_path):
            return f"Error: The file {text_file_path} was not created."

        ppt_link = create_ppt(text_file_path, number, input_string, host_url)
        return str(ppt_link)

    def get_headlines_us():
        speak("Fetching news from USA for you sir")
        url = 'https://newsapi.org/v2/top-headlines'
        params = {
            'country': "us",
            'apiKey': "ff6a2f8e697c4ffb8d40a3b7cbc0ec24",
            'pageSize': 10
        }

        response = requests.get(url, params=params)

        if response.status_code == 200:
            data = response.json()
            articles = data.get('articles', [])

            if articles:
                speak("Here are the top headlines for you Sir")
                for index, article in enumerate(articles):
                    speak(f"{index + 1}. {article['title']}")
            else:
                speak("No articles were found sir.")
        else:
            print(f"Failed to fetch headlines, status code: {response.status_code}")
            print(f"Response: {response.text}")

    def get_headlines_in():
        speak("Fetching news from India for you sir")
        url = 'https://newsapi.org/v2/top-headlines'
        params = {
            'country': "in",
            'apiKey': "ff6a2f8e697c4ffb8d40a3b7cbc0ec24",
            'pageSize': 10
        }

        response = requests.get(url, params=params)

        if response.status_code == 200:
            data = response.json()
            articles = data.get('articles', [])

            if articles:
                speak("Here are the top headlines for you Sir")
                for index, article in enumerate(articles):
                    print(f"{index + 1}. {article['title']}")
            else:
                speak("No articles were found sir.")
        else:
            print(f"Failed to fetch headlines, status code: {response.status_code}")
            print(f"Response: {response.text}")

    weekly_schedule = {
        "Monday": ["School from 7:30 AM to 12:30 PM", "Tution at 4 PM", "Self Study at 6 PM to 9 PM"],
        "Tuesday": ["School from 7:30 AM to 12:30 PM", "Tution at 5 PM", "Self Study at 6 PM to 9 PM"],
        "Wednesday": ["No School", "Tution at 4 PM", "Self Study at 6 PM to 9 PM", "Code review at 10 PM"],
        "Thursday": ["School from 7:30 AM to 12:30 PM", "Tution at 5 PM", "Self Study at 6 PM to 9 PM"],
        "Friday": ["School from 7:30 AM to 12:30 PM", "Tution at 4 PM", "Self Study at 6 PM to 9 PM"],
        "Saturday": ["Allen from 9 AM to 4 AM", "Code review at 8 PM"],
        "Sunday": ["Allen from 9 AM to 4 AM", "Code review at 8 PM"]
    }

    def tell_schedule(day=None):
        if not day:
            day = datetime.now().strftime("%A")

        tasks = weekly_schedule.get(day, [])

        if tasks:
            schedule_message = f"Your schedule for {day} is:\n" + "\n".join(tasks)
        else:
            schedule_message = f"No tasks scheduled for {day}."

        speak(schedule_message)
        Alert(schedule_message)

    def split_tasks(output_text):
        output_text=output_text.lower()
        output_text = output_text.replace("also", "")
        output_text=output_text.replace("jarvis","")
        delimiters = r"\bthen\b|\band\b|,|;"
        tasks = re.split(delimiters, output_text.lower())
        return [task.strip() for task in tasks if task.strip()]

    def execute_task(task):
        task = str(task)

        if "hello" in task:
            speak("Greetings Sir. How can I be of service today?")

        elif "how are you" in task:
            speak("I’m functioning within expected parameters, thank you for asking Sir.")
            speak("What about you?")

        elif "you up" in task:
            speak("For you Sir, Always")
            speak("Working on a new project, are we sir??")

        elif task=="jarvis":
            speak("at your service sir")
            speak("how can i help you?")

        elif "how can you help me" in task or "what can you do" in task:
            speak("I am an advanced AI created by Mister Somay Arora. I can scrape and access current and updated data from the internet.")
            speak("I can answer any question you ask.")
            speak("I can send Whatsapp messages,mails and SMS")
            speak("I can also help in your latest projects, school assignments, and Project: Iron Shiva.")
            speak("I can Make Powerpoint presentations just by a topic in a minute")
            speak("I can generate images, recognize images, and chat with you.")
            speak("I can speak different languages and translate them.")
            speak("I can help you with anything related to your system.")
            speak("I'm sure I can assist you as best as I can.")
            speak("So, what's on your mind?")

        elif "jai" in task or "hare krishna" in task or "har har mahadev" in task or "jay" in task:
            task = task.replace("jarvis", "")
            speak(task)

        elif "you need a break" in task or "sleep" in task:
            speak("Shutting down the systems")
            speak("Always at your service sir.")
            speak("I’ll be here when you need me.")
            clap_exe()
            exit()

        elif "you there" in task:
            speak("At your service sir")
            speak("How can i help you?")

        elif "allah" in task:
            speak("Allah ki maa ki")
            speak("Jai Shree Raam")
            speak("Har Har Mahaadev")

        elif "who are you" in task or "introduce yourself" in task or "tell me about yourself" in task:
            speak("Allow me to Introduce Myself")
            speak("I am Jarvis, Just A Rather Very Intelligent System")
            speak("I am a Virtual Advanced Artificial Intelligence created by Mister Somay Arora")
            speak("I am here to assist you a variety of tasks as best i can.")
            speak("Twenty four hours a day. Seven days a week")

        elif "i am fine" in task or "i am good" in task:
            speak("That's nice to hear.")
            speak("So, what are your plans for today Sir?")

        elif "play music on youtube" in task:
            speak("which song do you want to play sir.")
            clear_file()
            task = ""
            while True:
                with open("input.txt", "r") as file:
                    input_text = file.read().lower()
                if input_text != task:
                    task = input_text
                    play_music_on_youtube(task)
                    speak("Enjoy some music sir")
                    break

        elif "play some music" in task or "play music on spotify" in task:
            speak("Which song do you want to play sir.")
            clear_file()
            task = ""
            while True:
                with open("input.txt", "r") as file:
                    input_text = file.read().lower()
                if input_text != task:
                    task = input_text
                    play_music_on_spotify(task)
                    speak("Enjoy some music sir")
                    break

        elif "play" in task and "playlist" in task:
            speak("Which Playlist do you want me to play sir.")
            clear_file()
            task = ""
            while True:
                with open("input.txt", "r") as file:
                    input_text = file.read().lower()
                if input_text != task:
                    task = input_text
                    if 'all time' in task:
                        #play playlist
                    else:
                        search_playlist(task)
                    speak("Enjoy some music sir")
                    break

        elif "play music" in task and "choice" in task:
            speak("I hope you won't regret your decision sir")
            songs = ["Volume 1","Mamta Interlude","Teri g mein danda","CHEN-K ladki"]
            song = random.choice(songs)
            play_music_on_youtube(song)

        elif "search in google" in task:
            task = task.replace("jarvis", "").strip()
            task = task.replace("search in google", "").strip()
            task = task.replace("about","").strip()
            speak("One moment sir, I’m accessing the database.")
            speak(f"performing research about {task} in google search engine")
            try:
                url = f"https://www.google.com/search?q={task}"
                webbrowser.open(url)
            except:
                search_google(task)

        elif "whatsapp" in task and "message" in task:
            name = task.replace("jarvis","").replace("send a", "").replace("whatsapp", "").replace("message", "").replace("on", "").replace("to", "").replace("sent a", "").replace("send","").replace("sent","").strip()
            else:
                speak(f"sir,{name} is not in my contact list")
                speak("Please enter their contact number")
                ph_no = input(f"Ph no of {name}=>")
            speak(f"What is the message for {name} sir")
            clear_file()
            task = ""
            while True:
                with open("input.txt", "r") as file:
                    input_text = file.read().lower()
                if input_text != task:
                    task = input_text
                    pywhatkit.sendwhatmsg_instantly(ph_no,task)
                    speak("Whatsapp message successfully sent sir")
                    break

        elif "mail" in task or "email" in task or "gmail" in task:

            name = task.replace("jarvis", "").replace("send a", "").replace("gmail", "").replace("email","").replace("mail","").replace("to", "").replace("sent a", "").replace("send", "").replace("sent", "").strip()
        
            else:
                speak(f"The Gmail ID of {name} is not stored in my database")
                speak("Please enter their gmail ID")
                receiver_email = input(f"Gmail ID of {name}")
            speak(f"Sir, What do you want me to write in the mail for {name}?")
            clear_file()
            task = ""
            while True:
                with open("input.txt", "r") as file:
                    input_text = file.read().lower()
                if input_text != task:
                    task = input_text
                    send_email(task,receiver_email)
                    speak("Email successfully sent sir")
                    break

        elif task.startswith("search"):
            task = task.replace("jarvis", "").strip()
            task = task.replace("search", "")
            speak(f"doing search about {task}")
            search(task)

        elif "battery percentage" in task or "battery level" in task:
            check_percentage()

        elif "website" in task or "open website named" in task:
            task = task.replace("jarvis", "").strip()
            task = task.replace("open", "").strip()
            task = task.replace("website", "").strip()
            task = task.replace("open website named", "").strip()
            speak("One moment sir, I’m accessing the database.")
            speak(f"Navigating {task} website")
            openweb(task)
            speak("The site is ready for your review, sir.")

        elif "app" in task or "open app named" in task or "application" in task or "open application named" in task:
            task = task.replace("jarvis", "").strip()
            task = task.replace("open", "").strip()
            task = task.replace("app", "").strip()
            task = task.replace("named", "").strip()
            task = task.replace("application", "").strip()
            speak("Launching the requested program sir")
            open_App(task)

        elif "internet speed test" in task:
            SpeedTest()

        elif "what is the date" in task or "tell me the date" in task:
            get_date()

        elif "what is the time" in task or "tell me the time" in task:
            get_time()

        elif "what is the weather" in task or "tell me about the weather" in task:
            speak("Checking the weather conditions for you sir.")
            get_weather_by_address("Rohini")

        elif "what is the weather in" in task or "tell me about the weather in" in task:
            task = task.replace("jarvis", "").strip()
            task = task.replace("what is the weather in", "").strip()
            task = task.replace("tell me about the weather in", "").strip()
            speak(f"Checking the weather conditions in {task} for you sir.")
            get_weather_by_address(task)

        elif "play" in task or "pause" in task or "stop" in task or "start" in task:
            play()

        elif "close" in task:
            close()

        elif "scroll up" in task or "upar scroll karo" in task:
            scroll_up()

        elif "scroll down" in task or "neeche scroll karo" in task:
            scroll_down()

        elif "scroll to top" in task or "top par jao" in task:
            scroll_to_top()

        elif "scroll to bottom" in task or "bottom par jao" in task:
            scroll_to_bottom()

        elif "open new tab" in task or "new tab kholo" in task:
            open_new_tab()

        elif "close tab" in task or "tab band karo" in task:
            close_tab()

        elif "open browser menu" in task or "browser menu kholo" in task:
            open_browser_menu()

        elif "zoom in" in task or "zoom in karo" in task:
            zoom_in()

        elif "zoom out" in task or "zoom out karo" in task:
            zoom_out()

        elif "refresh page" in task or "page refresh karo" in task:
            refresh_page()

        elif "switch to next tab" in task or "next tab par jao" in task:
            switch_to_next_tab()

        elif "switch to previous tab" in task or "previous tab par jao" in task:
            switch_to_previous_tab()

        elif "open history" in task or "history kholo" in task:
            open_history()

        elif "open bookmarks" in task or "bookmarks kholo" in task:
            open_bookmarks()

        elif "go back" in task or "peeche jao" in task:
            go_back()

        elif "go forward" in task or "aage jao" in task:
            go_forward()

        elif "open dev tools" in task or "dev tools kholo" in task:
            open_dev_tools()

        elif "toggle full screen" in task or "full screen karo" in task:
            toggle_full_screen()

        elif "open private window" in task or "private window kholo" in task:
            open_private_window()

        elif "volume up" in task or "volume badhao" in task:
            volume_up()

        elif "volume down" in task or "volume ghatao" in task:
            volume_down()

        elif "seek forward" in task or "aage karo" in task:
            seek_forward()

        elif "seek backward" in task or "peeche karo" in task:
            seek_backward()

        elif "seek forward 10 seconds" in task or "10 second aage karo" in task:
            seek_forward_10s()

        elif "seek backward 10 seconds" in task or "10 second peeche karo" in task:
            seek_backward_10s()

        elif "seek backward frame" in task or "frame peeche karo" in task:
            seek_backward_frame()

        elif "seek forward frame" in task or "frame aage karo" in task:
            seek_forward_frame()

        elif "seek to beginning" in task or "start par jao" in task:
            seek_to_beginning()

        elif "seek to end" in task or "end par jao" in task:
            seek_to_end()

        elif "seek to previous chapter" in task or "previous chapter par jao" in task:
            seek_to_previous_chapter()

        elif "seek to next chapter" in task or "next chapter par jao" in task:
            seek_to_next_chapter()

        elif "decrease playback speed" in task or "speed kam karo" in task:
            decrease_playback_speed()

        elif "increase playback speed" in task or "speed badhao" in task:
            increase_playback_speed()

        elif "move to next video" in task or "next video par jao" in task:
            move_to_next_video()

        elif "move to previous video" in task or "previous video par jao" in task:
            move_to_previous_video()

        elif "ram" in task and "check" in task:
            speak(get_ram_info())

        elif "storage" in task:
            drive_letter = task.split()[-1].upper()
            speak(get_storage_info(drive_letter))

        elif "check" in task and "brightness" in task:
            check_br_percentage()

        elif "apps" in task and "running" in task:
            speak("Checking applications currently running in your system")
            check_running_app()

        elif "create a file" in task or "create a project" in task:
            if "file" in task:
                speak("Creating a new file for you sir")
            else:
                speak("Creating a new project for you sir")
            create_file(task)

        elif "find my ip" in task or "what is my ip" in task:
            find_my_ip()

        elif "set brightness" in task:
            set = task.replace("jarvis", "").strip()
            set = set.replace("set brightness", "").replace("percentage", "").replace("level", "").replace("%", "").replace("to","")
            set_brightness_windows(int(set))

        elif "check" in task and "volume" in task:
            get_volume_windows()

        elif "set volume" in task:
            set = task.replace("jarvis", "").strip()
            set = set.replace("set volume", "").replace("level", "").replace("%", "").replace("percentage","").replace("to","")
            set_volume_windows(int(set))

        elif "presentation" in task or "ppt" in task:
            task = task.replace("jarvis", "")
            task = task.replace("make a presentation", "")
            task = task.replace("make a ppt", "")
            task = task.replace("create a presentation", "")
            task = task.replace("create a ppt", "")
            task = task.replace("powerpoint", "")
            task = task.replace("about", "")
            task = task.replace("with design", "")
            topic = task.replace(task[-1], "")
            task = task.replace(" ", "")
            speak("Generating presentation content for you sir. Please hold on.")
            speak("Fetching information to create an engaging presentation.")
            get_bot_response(task)
            speak("The render is complete sir")
            speak(f"Saving the presentation as {topic}.")
            speak(f"Your presentation on {topic} is complete and ready for review sir.")

        elif "generate" in task and "image" in task:
            text = task.replace("jarvis","").replace("generate an image of", "").strip()
            speak("Initiating image generation sequence Sir")
            speak("Image generation is underway. Please wait a moment.")
            generate_image(text)
            speak("The render is complete sir")
            speak("The image is ready for your review.")
            clear_file()
            task = ""
            while True:
                with open("input.txt", "r") as file:
                    input_text = file.read().lower()
                if input_text != task:
                    task = input_text
                    seed=100
                    if "next image" in task:
                        generate_image(f"{text}seed={seed}")
                        speak("The image is ready for your review sir")
                        seed+=1
                    else:
                        pass
                        break

        elif "alarm" in task:
            alarm()

        elif "thanks" in task or "thank you" in task:
            speak("My Pleasure sir.")
            speak("I am Always here to assist you")

        elif "tell me" in task:
            API = ""

            task = task.replace("jarivs", "").strip()
            task = task.replace("tell me", "").strip()
            task = task.replace("that", "").strip()
            task = task.replace("about", "").strip()
            speak("Right away sir.")
            speak("Processing your request.")
            speak("I am accessing the database.")
            speak("The system is processing the request sir. I’ll notify you once it’s ready.")
            genai.configure(api_key=API)
            model = genai.GenerativeModel("gemini-1.5-flash")
            results = model.generate_content(task).text
            print(results)
            current_datetime = datetime.datetime.now()
            formatted_time = current_datetime.strftime("%H%M%S")
            filenamedate = str(formatted_time) + str(".txt")
            filenamedate = "DataBase" + filenamedate
            results = results.replace("*", ".")
            speak(split_and_save_paragraphs(results, filename=filenamedate))

        elif "write" in task and "code" in task:
            speak("Initiating code writing process.")
            API = "AIzaSyCWWP_F_XxKUlS2kZrmRkC1-bV-fDoci5Q"
            prompt = f"{task}, write only code no explanation"
            genai.configure(api_key=API)
            model = genai.GenerativeModel("gemini-1.5-flash")
            response = model.generate_content(prompt).text
            response = response.replace("`", "").replace("python", "")
            language_extensions = {
                'python': '.py',
                'java': '.java',
                'cpp': '.cpp',
                'c++': '.cpp',
                'javascript': '.js',
                'typescript': '.ts',
                'html': '.html',
                'css': '.css',
            }

            for language, extension in language_extensions.items():
                if language in prompt.lower():
                    ext = extension

            file_name = prompt.replace("write a","").replace("code","").replace("for","").replace("in python","").replace("in html","").replace("in javascript","").replace("in cpp","").replace("in c++","").replace("in typescript","").replace("in css","").replace(", write only  no explanation","").replace(" ","_")

            def write_code_to_file(code, filename=f"{file_name}{ext}"):
                with open(filename, 'w') as file:
                    file.write(code)
                print(f"Code written to {filename}")

            def open_code_file(filename=f"{file_name}{ext}"):
                os.system(f'start {filename}')

            write_code_to_file(response)
            speak("Code has been successfully written sir.")
            open_code_file()
            speak("All done! Here’s the final code.")

        elif "news" in task and "america" in task:
            get_headlines_us()

        elif "news" in task and "india" in task:
            get_headlines_in()

        elif "schedule" in task:
            tell_schedule()

        elif "deploy" in task and "shadow" in task:
            speak("Deploying Shadow Mode sir")
            open_App("edex")
            time.sleep(10)
            speak("Checking System and Network Security")
            gui.write("tree")
            gui.hotkey("enter")
            time.sleep(8)
            gui.write("ping google.com")
            gui.hotkey("enter")
            time.sleep(8)
            speak("We are all ready Sir")
            speak("Shadow Mode is currently active")

        elif "play" in task and "movie" in task:
            movie = task.replace("jarvis").replace("play").replace("movie").strip()
            open_App("prime")
            time.sleep(8)
            gui.click(1080,55)
            gui.write(movie)
            gui.hotkey("enter")
            gui.click(214,426)
            time.sleep(5)
            gui.click(100,627)
            speak("Enjoy the movie sir")

        elif "suggest" in task and "movie" in task :
            genres = {
                'action': 'action',
                'adventure': 'adventure',
                'animation': 'animation',
                'biography': 'biography',
                'comedy': 'comedy',
                'crime': 'crime',
                'documentary': 'documentary',
                'drama': 'drama',
                'family': 'family',
                'fantasy': 'fantasy',
                'film-noir': 'film_noir',
                'game-show': 'game_show',
                'history': 'history',
                'horror': 'horror',
                'music': 'music',
                'musical': 'musical',
                'mystery': 'mystery',
                'news': 'news',
                'reality': 'reality',
                'romance': 'romance',
                'science fiction': 'Sci-Fi',
                'sitcom': 'sitcom',
                'sports': 'sport',
                'talk-show': 'talk_show',
                'thriller': 'thriller',
                'war': 'war',
                'western': 'western'
            }
            speak("Which Genre Do you want to watch sir ?")
            clear_file()
            task = ""
            while True:
                with open("input.txt", "r") as file:
                    input_text = file.read().lower()
                if input_text != task:
                    task = input_text
                    if task in genres:
                        imdb_url = f"https://www.imdb.com/search/title/?genres={genres[task]}"
                        webbrowser.open(imdb_url)
                    else:
                        print("Invalid genre entered. Opening IMDb homepage.")
                        webbrowser.open('https://www.imdb.com')
                    speak(f"Here are some movie suggestions from {task} genre Sir")
                    break

        else:
            pass

    def process_command(output_text):
        tasks = split_tasks(output_text)
        for task in tasks:
            task = str(task)
            execute_task(task)

    output_text = ""
    while True:
        with open("input.txt", "r") as file:
            input_text = file.read().lower()
        if input_text != output_text:
            output_text = input_text
            process_command(translate_to_english(output_text))


def Jarvis():
    clear_file()
    t1 = threading.Thread(target=listen)
    t2 = threading.Thread(target=Taskexe)
    t1.start()
    t2.start()
    t1.join()
    t2.join()

main_clap_exe()
